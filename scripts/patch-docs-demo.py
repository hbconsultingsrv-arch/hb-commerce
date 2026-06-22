#!/usr/bin/env python3
"""Ajoute les sections demo aux documents Office restaurés sans les casser."""

from __future__ import annotations

from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
DOCX = DOCS / "cahier-des-charges-hb-commerce.docx"
PPTX = DOCS / "presentation-hb-commerce.pptx"
SITE = "https://hbconsultingsrv-arch.github.io/hb-commerce"

DEMO_ACCOUNTS = [
    ("super@hbcommerce.demo", "Super root", "super-root.html"),
    ("admin@hbcommerce.demo", "Admin", "admin.html"),
    ("agent.martin@hbcommerce.demo", "Agent commercial", "admin.html"),
    ("contact@restaurant-paris.demo", "Client — Le Jasmin", "compte.html"),
    ("stock@fiafi-tunisie.demo", "Fournisseur FIAFI", "supplier.html"),
]

EXAMPLE_PAGES = [
    "index.html — Vitrine",
    "produits.html — Catalogue",
    "compte.html — Espace client",
    "admin.html — Dashboard admin",
    "docs/exemples/index.html — Hub demo",
]


def esc(text: str) -> str:
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def docx_paragraph(text: str = "", style: str | None = None, bold: bool = False) -> str:
    props = f'<w:pStyle w:val="{style}"/>' if style else ""
    rpr = "<w:b/>" if bold else ""
    return (
        f"<w:p><w:pPr>{props}</w:pPr>"
        f'<w:r><w:rPr>{rpr}</w:rPr><w:t xml:space="preserve">{esc(text)}</w:t></w:r></w:p>'
    )


def patch_docx() -> None:
    insert = [
        docx_paragraph(""),
        docx_paragraph("12. Environnement de démonstration", "Heading1"),
        docx_paragraph("Site demo : " + SITE),
        docx_paragraph("Mot de passe commun : Test1234!"),
        docx_paragraph("Script SQL : supabase/seed-demo-data.sql"),
        docx_paragraph("Hub interactif : docs/exemples/index.html"),
        docx_paragraph(""),
        docx_paragraph("Comptes de démonstration", "Heading2"),
        docx_paragraph("E-mail | Rôle | Page principale", bold=True),
    ]
    for email, role, page in DEMO_ACCOUNTS:
        insert.append(docx_paragraph(f"{email} | {role} | {page}"))
    insert.extend([
        docx_paragraph(""),
        docx_paragraph("13. Pages d'exemple", "Heading1"),
    ])
    for line in EXAMPLE_PAGES:
        insert.append(docx_paragraph(f"{SITE}/{line.split(' — ')[0]} — {line.split(' — ', 1)[1]}"))
    block = "".join(insert)

    tmp = DOCS / "_docx_patch"
    if tmp.exists():
        import shutil
        shutil.rmtree(tmp)
    tmp.mkdir()

    with ZipFile(DOCX, "r") as zin, ZipFile(tmp / "out.docx", "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "word/document.xml":
                xml = data.decode("utf-8")
                marker = "<w:sectPr>"
                if marker not in xml:
                    raise RuntimeError("Structure Word inattendue")
                xml = xml.replace(marker, block + marker, 1)
                data = xml.encode("utf-8")
            zout.writestr(item, data)

    backup = DOCS / "cahier-des-charges-hb-commerce.backup.docx"
    if not backup.exists():
        backup.write_bytes(DOCX.read_bytes())
    (tmp / "out.docx").replace(DOCX)


def _set_title(slide, title: str) -> None:
    if not slide.shapes.title:
        return
    slide.shapes.title.text = title
    for paragraph in slide.shapes.title.text_frame.paragraphs:
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0x1A, 0x5C, 0x4A)


def _add_bullets(slide, lines: list[str]) -> None:
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def patch_pptx() -> None:
    prs = Presentation(str(PPTX))
    layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(layout)
    _set_title(slide, "Environnement de démonstration")
    _add_bullets(slide, [
        f"Site : {SITE}",
        "Mot de passe : Test1234! (tous les comptes)",
        "Script : supabase/seed-demo-data.sql",
        "Hub : docs/exemples/index.html",
        "super@ / admin@ / agent.martin@ / contact@restaurant-paris.demo",
    ])

    slide = prs.slides.add_slide(layout)
    _set_title(slide, "Pages d'exemple")
    _add_bullets(slide, [
        f"{SITE}/index.html — Vitrine",
        f"{SITE}/produits.html — Catalogue",
        f"{SITE}/compte.html — Espace client",
        f"{SITE}/admin.html — Dashboard admin",
        f"{SITE}/docs/exemples/index.html — Hub demo",
    ])

    backup = DOCS / "presentation-hb-commerce.backup.pptx"
    if not backup.exists():
        backup.write_bytes(PPTX.read_bytes())
    prs.save(str(PPTX))


def main() -> None:
    patch_docx()
    patch_pptx()
    with ZipFile(DOCX) as z:
        assert z.testzip() is None
    print("Documents patches avec sections demo.")


if __name__ == "__main__":
    main()
